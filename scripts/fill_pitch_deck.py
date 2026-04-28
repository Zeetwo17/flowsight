"""Fill the GSC 2026 prototype template with FlowSight content.

Reads the template, adds content text boxes below each section heading,
saves to artifacts/FlowSight_Pitch.pptx.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
TEMPLATE = Path.home() / "Downloads" / "[EXT] Solution Challenge 2026 - Prototype PPT Template (1).pptx"
OUT = REPO / "artifacts" / "FlowSight_Pitch.pptx"

# FlowSight brand
NAVY = RGBColor(0x1E, 0x29, 0x3A)        # primary text
ACCENT = RGBColor(0x32, 0x82, 0xFA)      # blue accent
RISK_RED = RGBColor(0xDC, 0x3C, 0x3C)
GOOD_GREEN = RGBColor(0x46, 0xC8, 0x78)
MUTED = RGBColor(0x55, 0x6B, 0x80)


def add_body(
    slide,
    content,
    top_in=1.55,
    left_in=0.4,
    width_in=9.3,
    height_in=4.05,
    head_pt=14,
    body_pt=11,
):
    """Add a body text box. `content` is a list of (text, level, bold) tuples."""
    tb = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = tb.text_frame
    tf.word_wrap = True

    first = True
    for entry in content:
        if isinstance(entry, str):
            text, level, bold = entry, 0, False
        else:
            text, level, bold = entry + (False,) * (3 - len(entry))
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = ""
        p.level = level
        run = p.add_run()
        run.text = text
        f = run.font
        f.name = "Calibri"
        f.size = Pt(head_pt if bold else body_pt)
        f.bold = bool(bold)
        f.color.rgb = ACCENT if bold else NAVY
    return tb


def trim_title_to(slide, keep_first_line: str) -> None:
    """Strip multi-paragraph placeholder titles down to the first line.

    The template's title boxes mix the heading with prompt sub-bullets
    ('Team name:', 'How different...?', etc.) - we replace the body
    with our own content, so the prompts become noisy duplicates.
    """
    ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        if sh.text_frame.paragraphs[0].text.strip() != keep_first_line.strip():
            continue
        tx_body = sh.text_frame._txBody
        paragraphs = tx_body.findall(f"{{{ns}}}p")
        for extra in paragraphs[1:]:
            tx_body.remove(extra)
        return


def add_callout(slide, label, value, left_in, top_in, width_in=2.0, height_in=1.2):
    """Big number + small label card."""
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = value
    r1.font.name = "Calibri"
    r1.font.size = Pt(34)
    r1.font.bold = True
    r1.font.color.rgb = ACCENT

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = label
    r2.font.name = "Calibri"
    r2.font.size = Pt(11)
    r2.font.bold = False
    r2.font.color.rgb = MUTED
    return box


def fill_team(slide):
    """Slide 2's title shape is mid-slide (template design) — fill it
    directly rather than overlay a body box that would collide with it."""
    target = None
    for sh in slide.shapes:
        if sh.has_text_frame and sh.text_frame.paragraphs[0].text.strip() == "Team Details":
            target = sh
            break
    if target is None:
        return
    tf = target.text_frame
    tf.clear()
    lines = [
        ("Team Details", True, 16),
        ("", False, 12),
        ("Team name:  FlowSight", True, 14),
        ("Team leader:  [your name here]", True, 14),
        ("", False, 12),
        ("Problem Statement", True, 14),
        ("Design a scalable system capable of continuously analyzing multifaceted", False, 12),
        ("transit data to preemptively detect and flag potential supply chain", False, 12),
        ("disruptions. Formulate dynamic mechanisms that instantly execute or", False, 12),
        ("recommend highly optimized route adjustments before localized", False, 12),
        ("bottlenecks cascade into broader delays.", False, 12),
    ]
    first = True
    for text, bold, sz in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Pt(sz)
        run.font.bold = bold
        run.font.color.rgb = ACCENT if bold else NAVY


def fill_brief(slide):
    add_body(slide, [
        ("FlowSight is a digital twin for supply chain ripple-effect mitigation.", 0, True),
        ("", 0, False),
        ("It ingests multifaceted transit data (news events, weather, traffic, historical", 0, False),
        ("patterns) into a single canonical event stream, models how disruptions cascade", 0, False),
        ("across a transit network, and recommends - or autonomously executes - risk-aware", 0, False),
        ("route adjustments before localized bottlenecks cascade into broader delays.", 0, False),
        ("", 0, False),
        ("Three layers, end-to-end:", 0, True),
        ("1.  Detection - spatial Hawkes / Neural Hawkes Process (cascade prediction)", 0, False),
        ("2.  Representation - Hidden Markov Model state per node (Normal / Stressed / Critical)", 0, False),
        ("3.  Decision - risk-aware Dijkstra + CVaR-PPO (route under tail risk)", 0, False),
        ("", 0, False),
        ("Every decision is explained in plain language by Gemini, so dispatchers see", 0, False),
        ("WHY a route changed - not just that it did.", 0, False),
    ], top_in=1.55, head_pt=14, body_pt=12)


def fill_opportunities(slide):
    trim_title_to(slide, "Opportunities")
    add_body(slide, [
        ("How is this different from existing systems?", 0, True),
        ("Conventional routing reacts to congestion AFTER it happens.", 0, False),
        ("FlowSight predicts the cascade and reroutes BEFORE failures spread.", 0, False),
        ("", 0, False),
        ("How does it solve the problem?", 0, True),
        ("Hawkes self-exciting point process models event propagation;", 0, False),
        ("HMM converts noisy intensities into actionable state labels;", 0, False),
        ("Dijkstra with CVaR penalty picks routes robust to tail risk -", 0, False),
        ("not just expected risk - critical for shipments with hard deadlines.", 0, False),
        ("", 0, False),
        ("USP - Unique Selling Proposition", 0, True),
        ("Research-backed end-to-end (Mei & Eisner 2017, Li et al. 2018,", 0, False),
        ("Schulman 2017, Rockafellar & Uryasev 2000, Ivanov & Dolgui 2020)", 0, False),
        ("Multi-source data fusion (GDELT news + weather + traffic + history)", 0, False),
        ("Gemini-powered plain-language reasoning for every reroute", 0, False),
        ("Sustainability built into the cost function: CO2 weighting (SDG 13)", 0, False),
        ("Dual mode: autonomous below confidence threshold, advisory above", 0, False),
    ], top_in=1.55, head_pt=14, body_pt=12)


def fill_features(slide):
    add_body(slide, [
        ("Cascade Detection", 0, True),
        ("Spatial Hawkes + optional Neural Hawkes (PyTorch CT-LSTM)", 0, False),
        ("", 0, False),
        ("State Estimation", 0, True),
        ("HMM labels each node Normal / Stressed / Critical", 0, False),
        ("", 0, False),
        ("Risk-Aware Routing", 0, True),
        ("Dijkstra with CVaR tail-risk objective + optional PPO policy", 0, False),
        ("", 0, False),
        ("Multi-Source Ingestion", 0, True),
        ("GDELT news + Open-Meteo weather + traffic + historical baselines", 0, False),
        ("", 0, False),
        ("Gemini Explanation Layer", 0, True),
        ("Plain-language justification per reroute (template fallback offline)", 0, False),
        ("", 0, False),
        ("Resilience Metrics + Live Map", 0, True),
        ("Resilience triangle (Bruneau 2003), live cascade animation,", 0, False),
        ("naive vs risk-aware route comparison side-by-side", 0, False),
    ], top_in=1.55, head_pt=14, body_pt=12)


def fill_process_flow(slide):
    add_body(slide, [
        ("Multi-source events", 0, True),
        ("GDELT news  ·  Open-Meteo weather  ·  traffic  ·  historical", 0, False),
        ("                   v", 0, False),
        ("Canonical Event Stream  -- (timestamp, target, kind, severity)", 0, True),
        ("                   v", 0, False),
        ("Hawkes Detector  -- intensity per node lambda_i(t)", 0, True),
        ("                   v", 0, False),
        ("HMM State Estimator  -- {Normal, Stressed, Critical} per node", 0, True),
        ("                   v", 0, False),
        ("Risk-Aware Router  -- travel_time + lambda * risk + CVaR + CO2", 0, True),
        ("                   v", 0, False),
        ("Decision Layer  -- auto-execute (high confidence) or advise (low)", 0, True),
        ("                   v", 0, False),
        ("Gemini Explainer  -- plain-language justification of the reroute", 0, True),
        ("                   v", 0, False),
        ("Dispatcher / Driver  -- map, KPIs, audit log on React PWA", 0, True),
    ], top_in=1.55, head_pt=13, body_pt=11)


def fill_wireframes(slide):
    add_body(slide, [
        ("Live working prototype - publicly accessible:", 0, True),
        ("https://flowsight-opal.vercel.app", 0, True),
        ("", 0, False),
        ("Three-column dashboard layout:", 0, True),
        ("Sidebar: scenario sliders + Run / Explain buttons", 0, False),
        ("Top-left map: real OSM tiles, color-coded nodes (green/yellow/red),", 0, False),
        ("grey route = naive shortest path, blue = risk-aware path", 0, False),
        ("Top-right: live cascade chart (Hawkes intensity over time)", 0, False),
        ("Middle-right: route comparison table (hops, time, risk, CO2)", 0, False),
        ("Bottom-left: disruption event log with source attribution", 0, False),
        ("Bottom-right: 'Why we rerouted' - Gemini explanation panel", 0, False),
        ("Top KPI bar: time saved, risk avoided, CO2 delta, shipments", 0, False),
        ("", 0, False),
        ("Installable as a PWA on iOS, Android, and desktop home screens.", 0, False),
    ], top_in=1.55, head_pt=14, body_pt=12)


def fill_architecture(slide):
    add_body(slide, [
        ("Browser (PWA)  -- Vercel Edge", 0, True),
        ("React 18 + TypeScript + Tailwind + Leaflet + Recharts", 0, False),
        ("Service Worker (offline-capable static shell)", 0, False),
        ("                |  HTTPS  /api/*  rewrite", 0, False),
        ("                v", 0, False),
        ("Backend (FastAPI on Render)  -- Python 3.12", 0, True),
        ("REST: /health  /graph  /simulate  /explain   ·   WebSocket: /stream", 0, False),
        ("                v", 0, False),
        ("In-process pipeline", 0, True),
        ("HawkesDetector  ->  StateEstimator  ->  RiskAwareRouter", 0, False),
        ("Optional: NeuralHawkesDetector  ·  DCRNNForecaster  ·  PPOAgent", 0, False),
        ("                v", 0, False),
        ("Reasoning  ->  Gemini API (gemini-2.5-flash)", 0, True),
        ("", 0, False),
        ("Optional GCP scale-out (Terraform-ready):", 0, True),
        ("Pub/Sub (events)  ·  Cloud Run  ·  Firestore (state)", 0, False),
        ("BigQuery (analytics)  ·  Secret Manager  ·  Firebase Hosting", 0, False),
    ], top_in=1.5, head_pt=13, body_pt=11)


def fill_technologies(slide):
    # Two-column layout: left = backend/ML, right = frontend/cloud
    add_body(slide, [
        ("Backend - Python 3.12", 0, True),
        ("FastAPI · NumPy · NetworkX · hmmlearn · PyTorch", 0, False),
        ("", 0, False),
        ("Machine Learning / Statistics", 0, True),
        ("Spatial Hawkes (cascade detection)", 0, False),
        ("Neural Hawkes - CT-LSTM (Mei & Eisner, NeurIPS 2017)", 0, False),
        ("DCRNN (Li et al., ICLR 2018)", 0, False),
        ("HMM (Rabiner 1989)", 0, False),
        ("PPO + CVaR (Schulman 2017, Rockafellar 2000)", 0, False),
        ("", 0, False),
        ("Data sources (free, public)", 0, True),
        ("GDELT 2.0  ·  Open-Meteo  ·  OpenStreetMap (osmnx)", 0, False),
    ], top_in=1.5, left_in=0.4, width_in=4.5, height_in=4.0,
       head_pt=13, body_pt=11)
    add_body(slide, [
        ("Google AI", 0, True),
        ("Gemini 2.5 Flash for plain-language reroute explanations", 0, False),
        ("(via google-generativeai SDK)", 0, False),
        ("", 0, False),
        ("Frontend - React 18 + TypeScript", 0, True),
        ("Vite · Tailwind CSS · Leaflet · Recharts · vite-plugin-pwa", 0, False),
        ("", 0, False),
        ("Hosting / Cloud", 0, True),
        ("Render (backend - free Docker tier)", 0, False),
        ("Vercel (frontend - free hobby tier)", 0, False),
        ("Optional GCP: Cloud Run, Pub/Sub, Firestore,", 0, False),
        ("BigQuery, Firebase Hosting (Terraform module included)", 0, False),
    ], top_in=1.5, left_in=5.0, width_in=4.5, height_in=4.0,
       head_pt=13, body_pt=11)


def fill_cost(slide):
    add_body(slide, [
        ("Cost to run the prototype: $0", 0, True),
        ("", 0, False),
        ("Render (backend) - free Docker tier", 0, True),
        ("750 hours/month, sleeps after 15 min idle, 30-90s cold start", 0, False),
        ("", 0, False),
        ("Vercel (frontend) - hobby tier", 0, True),
        ("Unlimited deploys, generous bandwidth, edge proxy included", 0, False),
        ("", 0, False),
        ("Gemini 2.5 Flash - free tier", 0, True),
        ("~1500 requests/day, ample for the demo window", 0, False),
        ("", 0, False),
        ("GDELT, Open-Meteo, OpenStreetMap, GitHub - free, no key", 0, True),
        ("", 0, False),
        ("Production scale-out (post-competition)", 0, True),
        ("GCP Cloud Run + Pub/Sub + BigQuery: $0-50/month at moderate", 0, False),
        ("traffic; well within the $300 free credit for the first 90 days.", 0, False),
        ("Terraform configs already in /infra/terraform.", 0, False),
    ], top_in=1.55, head_pt=14, body_pt=12)


def fill_mvp(slide):
    add_body(slide, [
        ("Live, deployed, public:", 0, True),
        ("Demo:    https://flowsight-opal.vercel.app", 0, False),
        ("API:     https://flowsight-api-ylyx.onrender.com", 0, False),
        ("Code:    https://github.com/Zeetwo17/flowsight", 0, False),
        ("", 0, False),
        ("What you can do in the live MVP:", 0, True),
        ("Adjust scenario parameters in the sidebar", 0, False),
        ("Click Run simulation - cascade animates across the corridor map", 0, False),
        ("Compare naive (grey) vs risk-aware (blue) routes side-by-side", 0, False),
        ("Click Explain reroute - Gemini generates a plain-language justification", 0, False),
        ("See KPIs (time saved, risk avoided, CO2 delta) update in real time", 0, False),
        ("", 0, False),
        ("Synthetic-scenario benchmark (Delhi-Mumbai 32-node corridor):", 0, True),
    ], top_in=1.55, height_in=2.7, head_pt=13, body_pt=11)
    # KPI callouts at the bottom
    add_callout(slide, "time saved", "168 min", 0.4, 4.4)
    add_callout(slide, "risk avoided", "63%", 2.6, 4.4)
    add_callout(slide, "explainability", "Gemini", 4.8, 4.4)
    add_callout(slide, "deploy cost", "$0", 7.0, 4.4)


def fill_future(slide):
    add_body(slide, [
        ("Real-data integration", 0, True),
        ("Live GDELT poller (India)  ·  Open-Meteo  ·  Google Maps Roads API", 0, False),
        ("", 0, False),
        ("PPO production training", 0, True),
        ("Train CVaR-PPO on historical disruption corpus; ship when", 0, False),
        ("worst-case return beats Dijkstra.", 0, False),
        ("", 0, False),
        ("Causal attribution", 0, True),
        ("Pearl-style: 'this delay was 70% weather, 20% cascade, 10% traffic'.", 0, False),
        ("", 0, False),
        ("Multi-shipment fleet optimization", 0, True),
        ("Joint routing under shared risk; minimise system-wide resilience triangle.", 0, False),
        ("", 0, False),
        ("Production touchpoints", 0, True),
        ("WhatsApp/SMS via Twilio  ·  driver PWA  ·  Cloud Run + Pub/Sub + BigQuery", 0, False),
        ("", 0, False),
        ("SDG alignment", 0, True),
        ("9 Infrastructure  ·  11 Sustainable Cities  ·  12 Responsible Consumption  ·  13 Climate", 0, False),
    ], top_in=1.55, height_in=4.0, head_pt=13, body_pt=11)


def fill_links(slide):
    trim_title_to(slide, "Provide links to your:")
    add_body(slide, [
        ("GitHub Public Repository", 0, True),
        ("https://github.com/Zeetwo17/flowsight", 0, False),
        ("", 0, False),
        ("Demo Video Link (3 minutes)", 0, True),
        ("[to be added before final submission]", 0, False),
        ("", 0, False),
        ("MVP Link", 0, True),
        ("https://flowsight-opal.vercel.app", 0, False),
        ("", 0, False),
        ("Working Prototype Link", 0, True),
        ("https://flowsight-opal.vercel.app", 0, False),
        ("", 0, False),
        ("API endpoints (FastAPI backend on Render)", 0, True),
        ("Base URL: https://flowsight-api-ylyx.onrender.com", 0, False),
        ("/health  ·  /graph  ·  /simulate  ·  /explain  ·  /stream (WebSocket)", 0, False),
    ], top_in=1.55, head_pt=14, body_pt=12)


def main() -> int:
    if not TEMPLATE.exists():
        print(f"Template not found: {TEMPLATE}")
        return 1
    OUT.parent.mkdir(parents=True, exist_ok=True)

    p = Presentation(str(TEMPLATE))
    slides = list(p.slides)
    fillers = {
        2: fill_team,
        3: fill_brief,
        4: fill_opportunities,
        5: fill_features,
        6: fill_process_flow,
        7: fill_wireframes,
        8: fill_architecture,
        9: fill_technologies,
        10: fill_cost,
        11: fill_mvp,
        12: fill_future,
        13: fill_links,
    }
    for idx, fn in fillers.items():
        if idx <= len(slides):
            fn(slides[idx - 1])
            print(f"  filled slide {idx}")

    p.save(str(OUT))
    print(f"\nWrote {OUT}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
